"""Generate PowerPoint presentation for DataServer In-House project proposal.

Usage:
    uv run python scripts/generate_presentation.py
"""

import sys
from pathlib import Path

sys.stdout.reconfigure(encoding="utf-8")
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Color scheme
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF2, 0xF3, 0xF4)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)


def add_background(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a colored title bar at the top of a slide."""
    # Title bar background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = LIGHT_BLUE


def add_body_text(slide, text_lines, left=0.5, top=1.5, width=9, size=16, bold_first=False):
    """Add body text with bullet points."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
        if bold_first and i == 0:
            p.font.bold = True


def add_info_box(slide, title, items, left, top, width=4.2, height=3.5, color=ACCENT_BLUE):
    """Add a styled info box with title and bullet items."""
    # Box background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    # Title
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.1), Inches(width - 0.4), Inches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color

    # Items
    txBox2 = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.55), Inches(width - 0.4), Inches(height - 0.7)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)


def add_table(slide, headers, rows, left=0.5, top=1.8, col_widths=None):
    """Add a styled table to the slide."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    width = Inches(9) if col_widths is None else sum(Inches(w) for w in col_widths)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols, Inches(left), Inches(top), width, Inches(0.4 * num_rows)
    )
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GRAY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 1: Title =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_background(slide, DARK_BLUE)

    # Main title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "DataServer In-House"
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "電力先物データ分析基盤 構築提案"
    p2.font.size = Pt(28)
    p2.font.color.rgb = LIGHT_BLUE
    p2.alignment = PP_ALIGN.CENTER

    # Decorative line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3), Inches(4.2), Inches(4), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    # Date
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "2026年3月"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
    p3.alignment = PP_ALIGN.CENTER

    # ===== SLIDE 2: Background & Challenges =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "背景・課題", "なぜ電力先物データの蓄積・分析が必要か")

    add_body_text(slide, [
        "■ 現状の課題",
        "  • 電力先物市場データは日次で公開されるが、体系的な蓄積・分析基盤がない",
        "  • 手動でのデータ収集は非効率で、過去データの参照が困難",
        "  • 電力先物と関連市場（コモディティ・指数・金利）の相関分析ができていない",
        "",
        "■ 解決策",
        "  • JPXデリバティブ理論価格データを日次自動取得・蓄積するデータサーバーを構築",
        "  • 電力先物を中心に、関連市場データ（53種の原資産）を一元管理",
        "  • インハウスで低コスト運用し、将来的にAWSへスケールアウト",
    ], size=14)

    # ===== SLIDE 3: Data Overview =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "データ概要", "JPXデリバティブ理論価格データ")

    add_info_box(slide, "データソース", [
        "JPX日本取引所グループ公開データ",
        "ファイル形式: CSV (CP932)",
        "更新頻度: 毎営業日",
        "データ量: 約43,600行/日",
        "ファイル名: rb{YYYYMMDD}.csv",
    ], left=0.3, top=1.5, width=4.3, height=2.8)

    add_info_box(slide, "カラム構成（12列）", [
        "銘柄コード / 銘柄名称",
        "PUT/CAL（オプション種別）",
        "限月 / 権利行使価格",
        "清算価格 / 理論価格",
        "原資産価格 / ボラティリティ",
        "金利 / 残日数 / 原資産名称",
    ], left=5.0, top=1.5, width=4.7, height=2.8)

    # Data breakdown table
    add_table(slide,
        ["カテゴリ", "原資産数", "代表例"],
        [
            ["電力先物", "12種", "東・西エリア（ベース/日中/週間/年度）"],
            ["株価指数", "10種", "日経225, TOPIX, JPX日経400, NYダウ"],
            ["コモディティ", "15種", "金, 原油, LNG, ゴム, 大豆"],
            ["国債", "3種", "長期, 中期, 超長期"],
            ["その他", "13種", "日経平均VI, REIT, 無担保コール等"],
        ],
        top=4.8, col_widths=[1.8, 1.2, 6.0]
    )

    # ===== SLIDE 4: System Overview =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Phase 1: システム概要", "プロトタイプ構成")

    # Architecture boxes
    components = [
        ("Data/", "CSV Files\n(JPXデータ)", 0.3, 2.0, ORANGE),
        ("storage.py", "ストレージ\n抽象レイヤー", 2.5, 2.0, ACCENT_BLUE),
        ("csv_parser.py", "CSVパーサー\n(CP932対応)", 4.7, 2.0, ACCENT_BLUE),
        ("repository.py", "データアクセス\n抽象レイヤー", 6.9, 2.0, ACCENT_BLUE),
        ("SQLite DB", "market_data.db\n(ローカルDB)", 6.9, 4.2, GREEN),
    ]

    for title, desc, left, top, color in components:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.0), Inches(1.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        txBox = slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(top + 0.1), Inches(1.8), Inches(0.4)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        txBox2 = slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(top + 0.5), Inches(1.8), Inches(0.9)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = DARK_GRAY
        p2.alignment = PP_ALIGN.CENTER

    # Key point
    add_body_text(slide, [
        "■ 設計のポイント",
        "  • storage.py / repository.py による抽象化レイヤーにより、AWS移行時にコアロジックの変更不要",
        "  • 環境変数でローカル/AWS切り替え（DATA_STORAGE=local|s3, DB_BACKEND=sqlite|dynamodb）",
        "  • Python + uv によるモダンなプロジェクト管理",
    ], top=5.5, size=12)

    # ===== SLIDE 5: Database Design =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "データベース設計", "アクセスパターンベース設計（DynamoDB移行対応）")

    add_table(slide,
        ["アクセスパターン", "メソッド", "SQLite", "DynamoDB（将来）"],
        [
            ["日付の全データ取得", "get_by_date(date)", "WHERE trade_date=?", "PK=trade_date"],
            ["日付+原資産", "get_by_date_and_underlying()", "WHERE date AND name", "PK+SK prefix"],
            ["銘柄の時系列", "get_instrument_history()", "WHERE code ORDER BY date", "GSI1"],
            ["原資産一覧", "get_underlying_names()", "SELECT DISTINCT", "Scan / GSI"],
            ["データ一括投入", "bulk_insert(date, records)", "INSERT OR IGNORE", "BatchWriteItem"],
            ["取込ログ記録", "log_import()", "INSERT import_log", "別テーブル"],
        ],
        top=1.5, col_widths=[2.0, 2.8, 2.2, 2.0]
    )

    add_body_text(slide, [
        "■ UNIQUE制約: (trade_date, instrument_code) — 同一日の重複取り込み防止",
        "■ DynamoDB設計: PK=trade_date, SK=instrument_code + GSI(instrument_code, trade_date)",
    ], top=5.5, size=12)

    # ===== SLIDE 6: Tech Stack =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "技術スタック", "Phase 1 プロトタイプ")

    add_info_box(slide, "言語・ランタイム", [
        "Python 3.12",
        "uv パッケージマネージャー",
        "pandas（データ処理）",
        "python-pptx（資料自動生成）",
    ], left=0.3, top=1.5, width=2.9, height=2.5)

    add_info_box(slide, "データベース", [
        "SQLite（ローカル、ゼロコスト）",
        "WALモード（並列読み取り対応）",
        "将来: DynamoDB移行対応済み",
    ], left=3.5, top=1.5, width=2.9, height=2.5)

    add_info_box(slide, "自動化・インフラ", [
        "GitHub Actions（日次スケジュール）",
        "個人GitHubアカウント（無料枠）",
        "将来: AWS Lambda + EventBridge",
    ], left=6.7, top=1.5, width=3.0, height=2.5)

    add_body_text(slide, [
        "■ 全技術選定の基準: コスト最小化 + AWS移行容易性",
        "■ Phase 1は外部サービス依存なし、ローカルのみで完結",
    ], top=4.5, size=14)

    # ===== SLIDE 7: Demo Output =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "デモ: データ取り込み結果", "rb20260310.csv — 2026年3月10日分データ")

    add_info_box(slide, "取り込み結果サマリー", [
        "総レコード数: 43,637件",
        "原資産数: 54種",
        "先物（FUT）: 430件",
        "コールオプション（CAL）: 21,605件",
        "プットオプション（PUT）: 21,604件",
    ], left=0.3, top=1.5, width=4.3, height=2.5)

    add_info_box(slide, "電力先物データ", [
        "電力先物契約数: 124件",
        "東エリア: ベース/日中/週間/年度",
        "西エリア: ベース/日中/週間/年度",
        "限月: 2026年3月 〜 2028年2月",
        "清算価格: 11.67円 〜 25.47円/kWh",
    ], left=5.0, top=1.5, width=4.7, height=2.5)

    # Sample data
    add_table(slide,
        ["銘柄名称", "限月", "清算価格", "残日数", "エリア"],
        [
            ["FUT_EEB_260330", "202603", "19.11", "22", "東・ベース"],
            ["FUT_EEP_260330", "202603", "19.50", "22", "東・日中"],
            ["FUT_EWB_260330", "202603", "15.29", "22", "西・ベース"],
            ["FUT_EWP_260330", "202603", "15.65", "22", "西・日中"],
            ["FUT_EEB_270330", "202703", "17.43", "387", "東・ベース"],
            ["FUT_EWB_270330", "202703", "16.07", "387", "西・ベース"],
        ],
        top=4.5, col_widths=[2.5, 1.2, 1.5, 1.0, 2.8]
    )

    # ===== SLIDE 8: AWS Migration Roadmap =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "AWS移行ロードマップ", "3段階のスケーラビリティ戦略")

    stages = [
        ("ステージ1（現在）", "ローカル + GitHub", [
            "SQLite + ローカルファイル",
            "GitHub Actions日次実行",
            "コスト: ¥0/月",
        ], GREEN, 0.3),
        ("ステージ2（承認後）", "AWSサーバレス", [
            "DynamoDB（オンデマンド）",
            "Lambda + EventBridge",
            "S3データレイク",
            "コスト: <$10/月",
        ], ACCENT_BLUE, 3.5),
        ("ステージ3（拡張）", "本格運用", [
            "CloudFront + React",
            "API Gateway + Lambda",
            "Bedrock AI分析",
            "Cognito認証",
        ], ORANGE, 6.7),
    ]

    for title, subtitle, items, color, left in stages:
        # Stage box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.5), Inches(2.9), Inches(4.0)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.color.rgb = color
        shape.line.width = Pt(3)

        # Stage title
        txBox = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(1.6), Inches(2.6), Inches(0.4)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color

        # Stage subtitle
        txBox2 = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(2.0), Inches(2.6), Inches(0.3)
        )
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = DARK_GRAY

        # Items
        txBox3 = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(2.5), Inches(2.6), Inches(2.8)
        )
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf3.paragraphs[0]
            else:
                p = tf3.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(6)

    # Arrow indicators
    add_body_text(slide, [
        "■ 各ステージの移行はrepository.py / storage.pyのバックエンド実装追加のみ。コアロジック変更不要。",
    ], top=5.8, size=12)

    # ===== SLIDE 9: AWS Architecture =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "AWS サーバレス構成図", "ステージ2: DynamoDB + Lambda + S3 + EventBridge")

    aws_components = [
        ("EventBridge", "日次スケジューラ\n(cron)", 0.3, 2.0, ORANGE),
        ("Lambda", "スクレーピング\n+ データ処理", 2.7, 2.0, ACCENT_BLUE),
        ("S3", "CSVデータレイク\n(原本保管)", 5.1, 2.0, GREEN),
        ("DynamoDB", "デリバティブ\n価格データ", 7.5, 2.0, ACCENT_BLUE),
        ("Lambda", "API関数\n(データ照会)", 2.7, 4.2, ACCENT_BLUE),
        ("API Gateway", "REST API\n(将来拡張)", 5.1, 4.2, ORANGE),
        ("CloudFront", "Webダッシュボード\n(将来拡張)", 7.5, 4.2, GREEN),
    ]

    for title, desc, left, top, color in aws_components:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.1), Inches(1.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        txBox = slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(top + 0.1), Inches(1.9), Inches(0.35)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        txBox2 = slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(top + 0.5), Inches(1.9), Inches(0.9)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = DARK_GRAY
        p2.alignment = PP_ALIGN.CENTER

    add_body_text(slide, [
        "■ 全てサーバレス構成 — サーバー管理不要、従量課金で最小コスト",
    ], top=6.2, size=12)

    # ===== SLIDE 10: Cost Comparison =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "コスト比較", "段階ごとの運用コスト")

    add_table(slide,
        ["項目", "ステージ1\n(GitHub)", "ステージ2\n(AWS)", "備考"],
        [
            ["データベース", "¥0\n(SQLite)", "~$0.30/月\n(DynamoDB)", "オンデマンド\n43,600件/日書込"],
            ["ストレージ", "¥0\n(ローカル)", "~$0.01/月\n(S3)", "4MB/日\n≒1.2GB/年"],
            ["コンピュート", "¥0\n(GitHub Actions)", "~$0.01/月\n(Lambda)", "1日1回実行\n無料枠内"],
            ["スケジューラ", "¥0\n(GitHub Actions)", "~$0.01/月\n(EventBridge)", "1日1イベント"],
            ["月額合計", "¥0", "< $1/月", "年間 < $12"],
        ],
        top=1.5, col_widths=[2.0, 2.2, 2.2, 2.6]
    )

    add_body_text(slide, [
        "■ ステージ2でも月額$1以下で運用可能（ダッシュボード追加時でも<$10/月）",
        "■ DynamoDBオンデマンドモード: 使った分だけ課金、アイドル時ほぼゼロ",
    ], top=5.8, size=12)

    # ===== SLIDE 11: Schedule =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "スケジュール", "Phase 1〜3 タイムライン")

    add_table(slide,
        ["フェーズ", "期間", "成果物", "ステータス"],
        [
            ["Phase 1", "2026年3月", "CSV取込 + SQLiteデータベース", "完了"],
            ["Phase 1.5", "2026年4月", "Streamlitダッシュボード\n分析デモノートブック", "次ステップ"],
            ["Phase 2", "2026年4-5月", "GitHub Actions自動化\nJPXサイトスクレーピング", "計画中"],
            ["Phase 3\n(承認後)", "2026年下期", "AWS移行\nDynamoDB + Lambda", "将来計画"],
        ],
        top=1.5, col_widths=[1.5, 1.8, 3.5, 1.5]
    )

    add_body_text(slide, [
        "■ Phase 1は完了済み — 43,637件のデータ取り込みを実証",
        "■ Phase 1.5でプレゼン用の可視化デモを作成",
        "■ 社内承認後、Phase 3でAWSサーバレス移行を開始",
    ], top=4.8, size=14)

    # ===== SLIDE 12: Summary =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BLUE)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "まとめ・次のステップ"
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Decorative line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3), Inches(1.8), Inches(4), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    summary_items = [
        "✓ JPXデリバティブ理論価格データ（53種・43,600行/日）の日次蓄積基盤を構築",
        "✓ 電力先物12種を含む全市場データを一元管理",
        "✓ AWS DynamoDB + Lambda へのスケーラブルな移行パスを確保",
        "✓ Phase 1 コスト: ¥0 → AWS移行後も月額 <$10",
        "",
        "【承認事項】",
        "  1. Phase 1.5（ダッシュボード）の開発着手",
        "  2. GitHub Actions による日次自動取得の開始",
        "  3. AWS移行の時期・予算の検討",
    ]

    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(4.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(summary_items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE if not item.startswith("【") else LIGHT_BLUE
        p.space_after = Pt(8)
        if item.startswith("【"):
            p.font.bold = True

    # Save
    output_dir = Path(__file__).resolve().parent.parent / "docs"
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / "DataServer_InHouse_Proposal.pptx"
    prs.save(str(output_path))
    print(f"Presentation saved to: {output_path}")


if __name__ == "__main__":
    create_presentation()
