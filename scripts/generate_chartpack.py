"""Generate daily Morning Report PPTX chart pack.

Produces a professional 7-slide chart pack with matplotlib charts
embedded as images, suitable for morning distribution to trading/risk teams.

Usage:
    uv run python scripts/generate_chartpack.py
"""

import shutil
import sys
import tempfile
from pathlib import Path

sys.stdout.reconfigure(encoding="utf-8")
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from src.asset_taxonomy import (
    ASSET_TAXONOMY,
    CATEGORY_META,
    COMMODITY_CATEGORIES,
    get_display_name,
)
from src.commodity_query import (
    get_commodity_forward_curve,
    get_cross_commodity_snapshot,
)
from src.query import get_power_futures
from src.repository import get_repository

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
BG_DARK = "#0B0E14"
BG_CARD = "#141820"
TEXT_WHITE = "#F0F2F5"
TEXT_GRAY = "#8B929A"
BLUE = "#2563EB"
CYAN = "#0EA5E9"
GOLD = "#D4A843"
GREEN = "#10B981"
RED = "#EF4444"
PURPLE = "#A855F7"
ORANGE = "#F97316"

PPTX_BG = RGBColor(0x0B, 0x0E, 0x14)
PPTX_CARD = RGBColor(0x14, 0x18, 0x20)
PPTX_WHITE = RGBColor(0xF0, 0xF2, 0xF5)
PPTX_GRAY = RGBColor(0x8B, 0x92, 0x9A)
PPTX_BLUE = RGBColor(0x25, 0x63, 0xEB)
PPTX_CYAN = RGBColor(0x0E, 0xA5, 0xE9)
PPTX_GOLD = RGBColor(0xD4, 0xA8, 0x43)
PPTX_GREEN = RGBColor(0x10, 0xB9, 0x81)
PPTX_RED = RGBColor(0xEF, 0x44, 0x44)

# Power curve display config: (underlying_name, label, color)
POWER_CURVES = [
    ("電力(東・ベース)", "East Base", BLUE),
    ("電力(東・日中)", "East Peak", CYAN),
    ("電力(西・ベース)", "West Base", GOLD),
    ("電力(西・日中)", "West Peak", ORANGE),
]

CATEGORY_COLORS = {
    "energy": ORANGE,
    "metals": GOLD,
    "industrial": TEXT_GRAY,
    "agriculture": GREEN,
}

# Widescreen 16:9 dimensions in inches
SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5
CHART_DPI = 300


# ---------------------------------------------------------------------------
# Matplotlib theme helpers
# ---------------------------------------------------------------------------

def _apply_dark_theme():
    """Configure matplotlib for dark-themed charts."""
    plt.rcParams.update({
        "figure.facecolor": BG_DARK,
        "axes.facecolor": BG_CARD,
        "axes.edgecolor": TEXT_GRAY,
        "axes.labelcolor": TEXT_WHITE,
        "text.color": TEXT_WHITE,
        "xtick.color": TEXT_GRAY,
        "ytick.color": TEXT_GRAY,
        "grid.color": "#1E2330",
        "grid.alpha": 0.6,
        "legend.facecolor": BG_CARD,
        "legend.edgecolor": TEXT_GRAY,
        "legend.labelcolor": TEXT_WHITE,
        "font.size": 11,
        "axes.titlesize": 14,
        "axes.labelsize": 12,
        "figure.dpi": CHART_DPI,
        "savefig.dpi": CHART_DPI,
        "savefig.facecolor": BG_DARK,
        "savefig.bbox": "tight",
        "savefig.pad_inches": 0.15,
    })


# ---------------------------------------------------------------------------
# Data helpers
# ---------------------------------------------------------------------------

def _get_trade_dates(repo):
    """Return (latest_date, prev_date) from the import log."""
    log = repo.get_import_log()
    dates = sorted(
        {e["trade_date"] for e in log if e["status"] == "success"},
        reverse=True,
    )
    latest = dates[0] if dates else None
    prev = dates[1] if len(dates) > 1 else None
    return latest, prev


def _build_power_forward(repo, trade_date, underlying_name):
    """Return (months, prices) for a power underlying on a given date."""
    records = repo.get_by_date_and_underlying(trade_date, underlying_name)
    futures = [
        r for r in records
        if not r.get("put_call") and r.get("settlement_price") is not None
    ]
    # Keep only monthly contracts (6-digit contract_month)
    monthly = [
        f for f in futures
        if f.get("contract_month") and len(str(f["contract_month"])) == 6
    ]
    monthly.sort(key=lambda x: x["contract_month"])
    months = [str(f["contract_month"]) for f in monthly]
    prices = [f["settlement_price"] for f in monthly]
    return months, prices


# ---------------------------------------------------------------------------
# Chart rendering functions (each returns a PNG path)
# ---------------------------------------------------------------------------

def render_power_forward_curves(repo, trade_date, prev_date, tmpdir):
    """Slide 2: Power forward curves with prev-day overlay."""
    _apply_dark_theme()
    fig, ax = plt.subplots(figsize=(11, 5.5))

    for uname, label, color in POWER_CURVES:
        months, prices = _build_power_forward(repo, trade_date, uname)
        if months:
            ax.plot(months, prices, "-o", color=color, label=label,
                    linewidth=2, markersize=4, zorder=3)
        if prev_date:
            pm, pp = _build_power_forward(repo, prev_date, uname)
            if pm:
                ax.plot(pm, pp, "--", color=color, alpha=0.35,
                        linewidth=1.2, label=f"{label} (prev)")

    ax.set_title("Power Forward Curves  (JPY/kWh)", fontsize=15, fontweight="bold",
                 pad=12)
    ax.set_xlabel("Contract Month")
    ax.set_ylabel("Settlement Price (JPY/kWh)")
    ax.grid(True, linestyle="--", alpha=0.4)
    ax.legend(loc="upper right", fontsize=8, ncol=2, framealpha=0.7)
    # Rotate x labels for readability
    plt.xticks(rotation=45, ha="right", fontsize=8)
    ax.yaxis.set_major_formatter(mticker.FormatStrFormatter("%.2f"))

    path = Path(tmpdir) / "power_forward.png"
    fig.savefig(str(path))
    plt.close(fig)
    return str(path)


def render_metals_forward_curves(repo, trade_date, prev_date, tmpdir):
    """Slide 4: Gold & Platinum forward curves."""
    _apply_dark_theme()
    fig, axes = plt.subplots(1, 2, figsize=(11, 5), sharey=False)

    metals = [
        ("金", "Gold (JPY/g)", GOLD),
        ("白金", "Platinum (JPY/g)", "#C0C0C0"),
    ]

    for ax, (uname, title, color) in zip(axes, metals):
        curve = get_commodity_forward_curve(repo, trade_date, uname)
        months = [c["month"] for c in curve if c["settlement"] is not None]
        prices = [c["settlement"] for c in curve if c["settlement"] is not None]
        if months:
            ax.plot(months, prices, "-o", color=color, linewidth=2, markersize=5, zorder=3)
        if prev_date:
            pcurve = get_commodity_forward_curve(repo, prev_date, uname)
            pm = [c["month"] for c in pcurve if c["settlement"] is not None]
            pp = [c["settlement"] for c in pcurve if c["settlement"] is not None]
            if pm:
                ax.plot(pm, pp, "--", color=color, alpha=0.35, linewidth=1.2)
        ax.set_title(title, fontsize=12, fontweight="bold")
        ax.grid(True, linestyle="--", alpha=0.4)
        ax.tick_params(axis="x", rotation=45, labelsize=7)
        ax.yaxis.set_major_formatter(mticker.FormatStrFormatter("%.0f"))

    fig.suptitle("Precious Metals Forward Curves", fontsize=15, fontweight="bold", y=0.98)
    fig.tight_layout(rect=[0, 0, 1, 0.93])

    path = Path(tmpdir) / "metals_forward.png"
    fig.savefig(str(path))
    plt.close(fig)
    return str(path)


def render_energy_forward_curves(repo, trade_date, prev_date, tmpdir):
    """Slide 5: Dubai Crude & LNG forward curves."""
    _apply_dark_theme()
    fig, axes = plt.subplots(1, 2, figsize=(11, 5), sharey=False)

    energy = [
        ("ドバイ原油", "Dubai Crude Oil (JPY/kl)", RED),
        ("LNG(プラッツJKM)", "LNG Platts JKM (USD/MMBtu)", CYAN),
    ]

    for ax, (uname, title, color) in zip(axes, energy):
        curve = get_commodity_forward_curve(repo, trade_date, uname)
        months = [c["month"] for c in curve if c["settlement"] is not None]
        prices = [c["settlement"] for c in curve if c["settlement"] is not None]
        if months:
            ax.plot(months, prices, "-o", color=color, linewidth=2, markersize=5, zorder=3)
        if prev_date:
            pcurve = get_commodity_forward_curve(repo, prev_date, uname)
            pm = [c["month"] for c in pcurve if c["settlement"] is not None]
            pp = [c["settlement"] for c in pcurve if c["settlement"] is not None]
            if pm:
                ax.plot(pm, pp, "--", color=color, alpha=0.35, linewidth=1.2)
        ax.set_title(title, fontsize=12, fontweight="bold")
        ax.grid(True, linestyle="--", alpha=0.4)
        ax.tick_params(axis="x", rotation=45, labelsize=7)
        ax.yaxis.set_major_formatter(mticker.FormatStrFormatter("%.0f"))

    fig.suptitle("Energy Commodity Forward Curves", fontsize=15, fontweight="bold", y=0.98)
    fig.tight_layout(rect=[0, 0, 1, 0.93])

    path = Path(tmpdir) / "energy_forward.png"
    fig.savefig(str(path))
    plt.close(fig)
    return str(path)


def render_cross_market_bar(snapshot, tmpdir):
    """Slide 6: Horizontal bar chart of daily % changes."""
    _apply_dark_theme()

    items = [s for s in snapshot if s["change_pct"] is not None]
    items.sort(key=lambda x: abs(x["change_pct"]))

    if not items:
        # Render a placeholder
        fig, ax = plt.subplots(figsize=(11, 5.5))
        ax.text(0.5, 0.5, "No change data available", ha="center", va="center",
                fontsize=16, color=TEXT_GRAY, transform=ax.transAxes)
        path = Path(tmpdir) / "cross_market_bar.png"
        fig.savefig(str(path))
        plt.close(fig)
        return str(path)

    labels = [s["display_en"] for s in items]
    values = [s["change_pct"] for s in items]
    colors = [GREEN if v >= 0 else RED for v in values]

    fig_height = max(4, len(items) * 0.4 + 1.5)
    fig, ax = plt.subplots(figsize=(11, min(fig_height, 6.5)))

    bars = ax.barh(labels, values, color=colors, edgecolor="none", height=0.6)
    ax.set_xlabel("Daily Change (%)")
    ax.set_title("Cross-Market Daily Changes (%)", fontsize=15, fontweight="bold", pad=12)
    ax.axvline(0, color=TEXT_GRAY, linewidth=0.8)
    ax.grid(True, axis="x", linestyle="--", alpha=0.3)
    ax.tick_params(axis="y", labelsize=8)

    # Value labels on bars
    for bar, val in zip(bars, values):
        x_pos = bar.get_width()
        ha = "left" if x_pos >= 0 else "right"
        offset = 0.15 if x_pos >= 0 else -0.15
        ax.text(x_pos + offset, bar.get_y() + bar.get_height() / 2,
                f"{val:+.2f}%", va="center", ha=ha, fontsize=7, color=TEXT_WHITE)

    fig.tight_layout()
    path = Path(tmpdir) / "cross_market_bar.png"
    fig.savefig(str(path))
    plt.close(fig)
    return str(path)


# ---------------------------------------------------------------------------
# PPTX slide builders
# ---------------------------------------------------------------------------

def _set_bg(slide, color=PPTX_BG):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, text, left, top, width, height, size=16,
              color=PPTX_WHITE, bold=False, align=PP_ALIGN.LEFT):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def _add_header_bar(slide, title, subtitle=None):
    """Add a top header bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(SLIDE_WIDTH), Inches(0.95),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PPTX_CARD
    shape.line.fill.background()

    # Left accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(0.06), Inches(0.95),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = PPTX_BLUE
    accent.line.fill.background()

    _add_text(slide, title, 0.35, 0.12, 10, 0.5, size=24,
              color=PPTX_WHITE, bold=True)
    if subtitle:
        _add_text(slide, subtitle, 0.35, 0.55, 10, 0.35, size=11,
                  color=PPTX_GRAY)


def _add_image_slide(prs, title, subtitle, image_path):
    """Add a slide with header + centered chart image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_header_bar(slide, title, subtitle)

    # Center the image below the header
    img_top = 1.05
    img_max_w = SLIDE_WIDTH - 0.6
    img_max_h = SLIDE_HEIGHT - img_top - 0.3
    slide.shapes.add_picture(
        image_path, Inches(0.3), Inches(img_top),
        Inches(img_max_w), Inches(img_max_h),
    )
    return slide


def _build_table(slide, headers, rows, left, top, col_widths, row_height=0.35):
    """Build a styled table on a dark slide."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(Inches(w) for w in col_widths)

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(left), Inches(top),
        total_w, Inches(row_height * n_rows),
    )
    table = tbl_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    # Header
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PPTX_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(9)
        p.font.color.rgb = PPTX_WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data
    for ri, row_data in enumerate(rows):
        for ci, val in enumerate(row_data):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            # Alternating row colors
            cell.fill.solid()
            if ri % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(0x14, 0x18, 0x20)
            else:
                cell.fill.fore_color.rgb = RGBColor(0x1A, 0x1F, 0x2B)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)
            p.font.color.rgb = PPTX_WHITE
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table


def _color_change_cell(table, row_idx, col_idx, value):
    """Color a cell green/red based on numeric value sign."""
    if value is None:
        return
    cell = table.cell(row_idx, col_idx)
    p = cell.text_frame.paragraphs[0]
    if value > 0:
        p.font.color.rgb = PPTX_GREEN
    elif value < 0:
        p.font.color.rgb = PPTX_RED


# ---------------------------------------------------------------------------
# Slide creation
# ---------------------------------------------------------------------------

def slide_title(prs, trade_date):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    # Decorative top line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(SLIDE_WIDTH), Inches(0.08),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PPTX_BLUE
    shape.line.fill.background()

    _add_text(slide, "Daily Market Report", 1.5, 2.0, 10, 1.0,
              size=44, color=PPTX_WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, f"Trade Date: {trade_date}", 1.5, 3.3, 10, 0.5,
              size=22, color=PPTX_CYAN, align=PP_ALIGN.CENTER)

    # Divider
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.2),
        Inches(4.3), Inches(0.04),
    )
    div.fill.solid()
    div.fill.fore_color.rgb = PPTX_BLUE
    div.line.fill.background()

    _add_text(slide, "JPX Derivative & Commodity Analytics", 1.5, 4.5, 10, 0.5,
              size=16, color=PPTX_GRAY, align=PP_ALIGN.CENTER)

    # Footer branding
    _add_text(slide, "DataServer In-House", 0.3, 6.8, 5, 0.4,
              size=10, color=PPTX_GRAY)


def slide_commodity_table(prs, snapshot, trade_date):
    """Slide 3: Commodity snapshot table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_header_bar(slide, "Commodity Snapshot",
                    f"Front-month settlement prices as of {trade_date}")

    headers = ["Category", "Asset", "Contract", "Settlement", "Unit", "Change", "% Chg"]
    col_widths = [1.4, 2.5, 1.3, 1.5, 1.3, 1.3, 1.1]

    rows = []
    change_indices = []  # (row_idx_in_table, change_diff, change_pct)
    for s in snapshot:
        cat_label = CATEGORY_META.get(s["category"], {}).get("display_en", s["category"])
        diff_str = f"{s['change_diff']:+.2f}" if s["change_diff"] is not None else "-"
        pct_str = f"{s['change_pct']:+.2f}%" if s["change_pct"] is not None else "-"
        rows.append([
            cat_label,
            s["display_en"],
            str(s["contract_month"]),
            f"{s['settlement']:,.2f}",
            s["unit"],
            diff_str,
            pct_str,
        ])
        change_indices.append((len(rows), s.get("change_diff"), s.get("change_pct")))

    table = _build_table(slide, headers, rows, 0.4, 1.15, col_widths, row_height=0.32)

    # Color the change columns
    for row_1idx, diff, pct in change_indices:
        _color_change_cell(table, row_1idx, 5, diff)
        _color_change_cell(table, row_1idx, 6, pct)


def slide_power_movers(prs, repo, trade_date, prev_date):
    """Slide 7: Top 15 power futures movers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_header_bar(slide, "Top Power Futures Movers",
                    f"Largest absolute settlement price changes vs previous day ({trade_date})")

    current_futures = get_power_futures(repo, trade_date)
    current_map = {}
    for r in current_futures:
        if not r.get("put_call") and r.get("settlement_price") is not None:
            key = r["instrument_code"]
            current_map[key] = r

    prev_map = {}
    if prev_date:
        prev_futures = get_power_futures(repo, prev_date)
        for r in prev_futures:
            if not r.get("put_call") and r.get("settlement_price") is not None:
                prev_map[r["instrument_code"]] = r

    movers = []
    for code, cur in current_map.items():
        prev = prev_map.get(code)
        if prev and prev.get("settlement_price"):
            diff = cur["settlement_price"] - prev["settlement_price"]
            pct = diff / prev["settlement_price"] * 100 if prev["settlement_price"] else 0
            movers.append({
                "code": code,
                "name": cur.get("instrument_name", code),
                "underlying": cur.get("underlying_name", ""),
                "month": cur.get("contract_month", ""),
                "price": cur["settlement_price"],
                "prev": prev["settlement_price"],
                "diff": round(diff, 2),
                "pct": round(pct, 2),
            })

    movers.sort(key=lambda x: abs(x["diff"]), reverse=True)
    top = movers[:15]

    headers = ["Instrument", "Underlying", "Month", "Price", "Prev", "Change", "% Chg"]
    col_widths = [2.8, 2.5, 1.1, 1.2, 1.2, 1.2, 1.1]

    rows = []
    change_indices = []
    for m in top:
        rows.append([
            m["name"],
            get_display_name(m["underlying"]),
            str(m["month"]),
            f"{m['price']:.2f}",
            f"{m['prev']:.2f}",
            f"{m['diff']:+.2f}",
            f"{m['pct']:+.2f}%",
        ])
        change_indices.append((len(rows), m["diff"], m["pct"]))

    if not rows:
        _add_text(slide, "No previous-day data available for comparison.",
                  1, 3, 10, 1, size=16, color=PPTX_GRAY, align=PP_ALIGN.CENTER)
        return

    table = _build_table(slide, headers, rows, 0.3, 1.15, col_widths, row_height=0.32)
    for row_1idx, diff, pct in change_indices:
        _color_change_cell(table, row_1idx, 5, diff)
        _color_change_cell(table, row_1idx, 6, pct)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def generate_chartpack():
    """Main entry point: build and save the daily chart pack."""
    repo = get_repository()
    try:
        trade_date, prev_date = _get_trade_dates(repo)
        if not trade_date:
            print("ERROR: No data found in database. Run import first.")
            sys.exit(1)

        print(f"Generating chart pack for {trade_date} (prev: {prev_date})")

        # Fetch cross-commodity snapshot (used by slide 3 and slide 6)
        snapshot = get_cross_commodity_snapshot(repo, trade_date, prev_date)

        # Create temp directory for chart PNGs
        tmpdir = tempfile.mkdtemp(prefix="chartpack_")
        try:
            # Render all charts
            print("  Rendering power forward curves ...")
            power_img = render_power_forward_curves(repo, trade_date, prev_date, tmpdir)
            print("  Rendering metals forward curves ...")
            metals_img = render_metals_forward_curves(repo, trade_date, prev_date, tmpdir)
            print("  Rendering energy forward curves ...")
            energy_img = render_energy_forward_curves(repo, trade_date, prev_date, tmpdir)
            print("  Rendering cross-market bar chart ...")
            bar_img = render_cross_market_bar(snapshot, tmpdir)

            # Build PPTX
            prs = Presentation()
            prs.slide_width = Inches(SLIDE_WIDTH)
            prs.slide_height = Inches(SLIDE_HEIGHT)

            print("  Building slides ...")
            # Slide 1: Title
            slide_title(prs, trade_date)

            # Slide 2: Power Forward Curves
            _add_image_slide(prs, "Power Forward Curves",
                             f"East/West Base & Peak monthly contracts ({trade_date})",
                             power_img)

            # Slide 3: Commodity Snapshot Table
            slide_commodity_table(prs, snapshot, trade_date)

            # Slide 4: Metals Forward Curves
            _add_image_slide(prs, "Precious Metals Forward Curves",
                             f"Gold & Platinum ({trade_date})", metals_img)

            # Slide 5: Energy Forward Curves
            _add_image_slide(prs, "Energy Commodity Forward Curves",
                             f"Dubai Crude Oil & LNG Platts JKM ({trade_date})",
                             energy_img)

            # Slide 6: Cross-Market Daily Changes
            _add_image_slide(prs, "Cross-Market Daily Changes",
                             "All commodities sorted by absolute % change",
                             bar_img)

            # Slide 7: Top Power Movers
            slide_power_movers(prs, repo, trade_date, prev_date)

            # Save
            output_dir = Path(__file__).resolve().parent.parent / "docs" / "chartpack"
            output_dir.mkdir(parents=True, exist_ok=True)

            date_str = trade_date.replace("-", "")
            dated_path = output_dir / f"daily_chartpack_{date_str}.pptx"
            latest_path = output_dir / "latest.pptx"

            prs.save(str(dated_path))
            print(f"  Saved: {dated_path}")

            # Copy as latest.pptx (symlinks can be unreliable on Windows)
            shutil.copy2(str(dated_path), str(latest_path))
            print(f"  Copied: {latest_path}")

        finally:
            # Clean up temp chart images
            shutil.rmtree(tmpdir, ignore_errors=True)

    finally:
        repo.close()

    print("Done.")


if __name__ == "__main__":
    generate_chartpack()
